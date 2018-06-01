#The purpose of this application is to automatically import lyrics from an online source onto
# a created .pptx file
from bs4 import BeautifulSoup as soup
import bs4
import urllib.request
import requests
from pptx import Presentation
from pptx.shapes import *
import datetime
from pptx.util import *
from pptx.enum.text import *
#datetime.date.today()

class SlideShow:

	def __init__(self,title,artist):
		self._title=title
		self._artist=artist
		self._base='https://www.musixmatch.com'
		self._filename='Worship {}'.format(datetime.date.today().strftime('%m-%d'))
		self.headers = 	{'User-Agent': 'Mozilla/5.0 (Windows NT 6.1) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/41.0.2228.0 Safari/537.36'}
		self.grabSongData()

	def grabSongData(self):
		base_url='https://www.musixmatch.com/search/'
		songTitle=self._title.split(' ')
		songArtist=self._artist.split(' ')
		my_url=''
		songInfo=songTitle+songArtist

		for i in range(len(songInfo)-1):
			base_url+=songInfo[i]+'%20'

		my_url+=base_url+songInfo[-1]

		htmlPage= requests.get(my_url,headers=self.headers)
		htmlPage.raise_for_status()
		page_soup=soup(htmlPage.text,"html.parser")

		songUrl=page_soup.findAll("h2",{'class':'media-card-title'})[0]
		songUrl=songUrl.find('a',{'class':'title'})
		self._title=songUrl.span.string
		my_url=self._base+songUrl['href']

		htmlPage= requests.get(my_url,headers=self.headers)
		htmlPage.raise_for_status()
		page_soup=soup(htmlPage.text,"html.parser")

		songLyrics=page_soup.findAll('p',{'class':'mxm-lyrics__content'})
		songLyrics[0]=songLyrics[0].text
		songLyrics[1]=songLyrics[1].text
		self.splitLyrics(songLyrics)

	def splitLyrics(self,lyrics):
		refinedLyric=[]
		for verse in lyrics:
			verse=verse.replace('1','')
			verse=verse.replace('2','')
			verse=verse.replace('3','')
			verse=verse.replace('4','')			
			verse=verse.replace('BRIDGE','')
			refinedLyric+=(verse.split('\n\n'))
		for i in range(len(refinedLyric)):
			refinedLyric[i]=refinedLyric[i].strip()
		self._lyrics=refinedLyric


	def displayLyrics(self):
		for line in self._lyrics:
			print(line)
			print()

	def makeSlide(self):
		prs=Presentation()

		blank_slide_layout = prs.slide_layouts[5]
		for verse in self._lyrics:
			slide = prs.slides.add_slide(blank_slide_layout)
			shapes = slide.shapes
			left = top = width = height = Inches(1)
			txBox = slide.shapes.add_textbox(left+Inches(3.5), top+Inches(0.25), width, height)
			tf = txBox.text_frame

			title = slide.shapes.title
			title.text = self._title
			p = tf.add_paragraph()
			p.text = verse
			p.font.size = Pt(32)
			p.alignment = PP_ALIGN.CENTER

		prs.save(self._filename+'.pptx')
	def addSong(self,title,artist):
		prs=Presentation(self._filename+'.pptx')
		newSong=SlideShow(title,artist)

		blank_slide_layout = prs.slide_layouts[5]
		for verse in newSong._lyrics:
			slide = prs.slides.add_slide(blank_slide_layout)
			shapes = slide.shapes
			left = top = width = height = Inches(1)
			txBox = slide.shapes.add_textbox(left+Inches(3.5), top+Inches(0.25), width, height)
			tf = txBox.text_frame

			title = slide.shapes.title
			title.text = newSong._title
			p = tf.add_paragraph()
			p.text = verse
			p.font.size = Pt(32)
			p.alignment = PP_ALIGN.CENTER

		prs.save(self._filename+'.pptx')		

